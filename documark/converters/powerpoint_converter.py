"""
PowerPoint 转换器
"""

from pathlib import Path
from typing import Dict, Any, List

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .base import BaseConverter


class PowerPointConverter(BaseConverter):
    """
    PowerPoint演示文稿转换器
    
    将.ppt, .pptx文件转换为Markdown格式
    """
    
    CATEGORY = "演示"
    
    def convert(self, file_path: Path) -> str:
        """
        将PowerPoint转换为Markdown
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            Markdown格式的文本内容
        """
        if not PPTX_AVAILABLE:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
        
        prs = Presentation(file_path)
        content_parts = []
        
        # 添加文档标题
        content_parts.append(self._create_header(f"演示文稿: {file_path.stem}", 1))
        content_parts.append(f"**幻灯片数量**: {len(prs.slides)}\n\n")
        content_parts.append("---\n\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(self._process_slide(slide, slide_num))
        
        return "".join(content_parts)
    
    def _process_slide(self, slide, slide_num: int) -> str:
        """处理单个幻灯片"""
        content_parts = []
        
        content_parts.append(self._create_header(f"幻灯片 {slide_num}", 2))
        
        # 提取幻灯片中的文本
        slide_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                rows_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_data)
                
                if rows_data:
                    headers = rows_data[0]
                    data_rows = rows_data[1:] if len(rows_data) > 1 else []
                    slide_texts.append(self._create_table(headers, data_rows))
        
        if slide_texts:
            content_parts.append("\n\n".join(slide_texts))
        
        content_parts.append("\n\n---\n\n")
        
        return "".join(content_parts)
